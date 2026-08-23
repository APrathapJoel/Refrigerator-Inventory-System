import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    TOTAL_SLIDES = 11

    # --- Color Palette ---
    COLOR_BG = RGBColor(248, 250, 252)        # Slate 50
    COLOR_CARD_BG = RGBColor(255, 255, 255)   # Pure White
    COLOR_CARD_BORDER = RGBColor(203, 213, 225) # Slate 300
    COLOR_NAVY = RGBColor(15, 23, 42)          # Slate 900
    COLOR_PRIMARY = RGBColor(2, 132, 199)      # Sky 600
    COLOR_DARK_BLUE = RGBColor(30, 58, 138)    # Blue 900
    COLOR_ACCENT = RGBColor(14, 165, 233)      # Sky 500
    COLOR_EMERALD = RGBColor(16, 185, 129)     # Emerald 500
    COLOR_AMBER = RGBColor(245, 158, 11)       # Amber 500
    COLOR_ROSE = RGBColor(239, 68, 68)         # Rose 500
    COLOR_PURPLE = RGBColor(139, 92, 246)      # Purple 500
    COLOR_TEXT_MAIN = RGBColor(30, 41, 59)     # Slate 800
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_LINE = RGBColor(148, 163, 184)       # Slate 400

    def add_header(slide, title_text, category_text="PROTOTYPE SYSTEM SPECIFICATION"):
        # Top banner background bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_NAVY
        bar.line.fill.background()

        # Category Tracker
        tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.5), Inches(0.3))
        tf_cat = tb_cat.text_frame
        p_c = tf_cat.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(9.5)
        p_c.font.bold = True
        p_c.font.color.rgb = COLOR_ACCENT

        # Main Slide Title
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(11.5), Inches(0.65))
        tf_title = tb_title.text_frame
        p_t = tf_title.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE

    def add_footer(slide, slide_num):
        tb_f = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.35))
        tf_f = tb_f.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"Refrigerator & Pantry Inventory Subsystem (RIMS)   |   Prototype Engineering Specification   |   Slide {slide_num} of {TOTAL_SLIDES}"
        p_f.font.size = Pt(9)
        p_f.font.color.rgb = COLOR_TEXT_MUTED

    def create_card(slide, left, top, width, height, title, header_color=COLOR_PRIMARY):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER
        card.line.width = Pt(1.5)

        h_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.48))
        h_bar.fill.solid()
        h_bar.fill.fore_color.rgb = header_color
        h_bar.line.fill.background()
        
        tf_h = h_bar.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = f"  {title.upper()}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        return card

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_NAVY
    bg1.line.fill.background()

    tb_sub = s1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.9), Inches(0.5))
    p_sub = tb_sub.text_frame.paragraphs[0]
    p_sub.text = "PROTOTYPE SPECIFICATION & BLUEPRINT"
    p_sub.font.size = Pt(14)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_ACCENT

    tb_main = s1.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(10.9), Inches(2.2))
    p_main = tb_main.text_frame.paragraphs[0]
    p_main.text = "Refrigerator & Pantry\nInventory Management Subsystem"
    p_main.font.size = Pt(36)
    p_main.font.bold = True
    p_main.font.color.rgb = COLOR_WHITE

    p_desc = tb_main.text_frame.add_paragraph()
    p_desc.text = "Real-Time Stock Logging, 6 Food Categories, FEFO Expiration Alerts & Automated Reordering"
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = COLOR_TEXT_MUTED
    p_desc.space_before = Pt(14)

    badges = [
        ("Architecture", "Node / Express / SQLite", COLOR_PRIMARY),
        ("Frontend", "React 18 + Vite", COLOR_ACCENT),
        ("Storage Scope", "6 Food Categories", COLOR_EMERALD),
        ("Live Deployment", "Vercel Cloud", COLOR_PURPLE)
    ]
    for idx, (blabel, bval, bcol) in enumerate(badges):
        bx = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2 + idx * 2.75), Inches(5.1), Inches(2.55), Inches(1.0))
        bx.fill.solid()
        bx.fill.fore_color.rgb = RGBColor(30, 41, 59)
        bx.line.color.rgb = bcol
        bx.line.width = Pt(1.5)
        
        tf_bx = bx.text_frame
        p1 = tf_bx.paragraphs[0]
        p1.text = blabel.upper()
        p1.font.size = Pt(9)
        p1.font.bold = True
        p1.font.color.rgb = bcol
        
        p2 = tf_bx.add_paragraph()
        p2.text = bval
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(3)
    add_footer(s1, 1)

    # ==========================================
    # SLIDE 2: EXECUTIVE SUMMARY & PROTOTYPE SCOPE
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "1. Executive Summary & Prototype Scope")
    scope_cards = [
        ("1. Real-Time Pantry Tracking", "Tracks exact physical stock levels across 6 distinct food categories with instant balance updates.", COLOR_PRIMARY),
        ("2. Batch & Shelf-Life Management", "Maintains purchase and expiration dates per intake batch to enforce First-Expired, First-Out (FEFO) usage.", COLOR_DARK_BLUE),
        ("3. Instant Material Discovery", "Equipped with live category dropdown filters and keyword search inside modals to locate materials immediately.", COLOR_EMERALD),
        ("4. 5 Visual Telemetry Badges", "Real-time color-coded badges: Sufficient (🟢), Low Stock (🟡), Expiring Soon (🟠), Expired (🔴), and Out of Stock (🔴).", COLOR_AMBER),
        ("5. Proactive 3-Day Alerts", "Top alert banner flags any batch expiring within 3 days to prioritize usage and eliminate food waste.", COLOR_ROSE),
        ("6. Automated Reorder Engine", "Auto-calculates replenishment quantities when items fall below threshold, generating an instant Shopping List.", COLOR_PURPLE)
    ]
    for idx, (stitle, sdesc, scolor) in enumerate(scope_cards):
        r = idx // 3
        c = idx % 3
        card = create_card(s2, 0.8 + c * 3.95, 1.4 + r * 2.65, 3.8, 2.45, stitle, scolor)
        tb = s2.shapes.add_textbox(Inches(0.95 + c * 3.95), Inches(1.95 + r * 2.65), Inches(3.5), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sdesc
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.line_spacing = 1.3
    add_footer(s2, 2)

    # ==========================================
    # SLIDE 3: 6 FOOD CATEGORIES
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "2. Food Category Architecture & Storage Classification")
    cats = [
        ("1. Fruits", "Temp: 4°C - 8°C | High Humidity", "Apples, Bananas, Strawberries, Lemons", COLOR_ROSE),
        ("2. Vegetables", "Temp: 2°C - 6°C | High Humidity", "Carrots, Spinach, Roma Tomatoes, Garlic", COLOR_EMERALD),
        ("3. Dairy Products", "Temp: 1°C - 4°C | Controlled Chill", "Whole Milk, Cheddar Cheese, Butter, Heavy Cream", COLOR_ACCENT),
        ("4. Baking Products", "Temp: 15°C - 20°C | Dry & Sealed", "All-Purpose Flour, Granulated Sugar, Dry Yeast", COLOR_AMBER),
        ("5. Dessert Products", "Temp: 2°C - 5°C | Cool / Chilled", "Dark Chocolate Chips, Tart Shells, Maple Syrup", COLOR_PURPLE),
        ("6. Raw Materials & Other", "Temp: 15°C - 22°C | Ambient Pantry", "Olive Oil, Soy Sauce, Black Pepper, Sea Salt", COLOR_DARK_BLUE)
    ]
    for idx, (cname, ctemp, citems, ccolor) in enumerate(cats):
        r = idx // 3
        c = idx % 3
        create_card(s3, 0.8 + c * 3.95, 1.4 + r * 2.65, 3.8, 2.45, cname, ccolor)
        tb = s3.shapes.add_textbox(Inches(0.95 + c * 3.95), Inches(1.95 + r * 2.65), Inches(3.5), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = f"STORAGE PROFILE:\n{ctemp}"
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_NAVY
        p2 = tf.add_paragraph()
        p2.text = f"PROTOTYPE ITEMS:\n{citems}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_MAIN
        p2.space_before = Pt(8)
    add_footer(s3, 3)

    # ==========================================
    # SLIDE 4: OBJECT-ORIENTED UML CLASS DIAGRAM
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "3. Object-Oriented UML Class Diagram")

    # Function to draw a 3-compartment UML Class Box
    def draw_uml_class(slide, left, top, width, height, class_name, stereotype, attributes, methods, header_color=COLOR_NAVY):
        # Outer container box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD_BG
        box.line.color.rgb = header_color
        box.line.width = Pt(1.5)

        # 1. Header Compartment (Class Name & Stereotype)
        h_height = 0.58
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(h_height))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = header_color
        hdr.line.fill.background()
        tf_h = hdr.text_frame
        tf_h.word_wrap = True
        p_st = tf_h.paragraphs[0]
        p_st.text = f"«{stereotype}»"
        p_st.font.size = Pt(8)
        p_st.font.bold = True
        p_st.font.color.rgb = COLOR_ACCENT
        p_st.alignment = PP_ALIGN.CENTER
        
        p_cn = tf_h.add_paragraph()
        p_cn.text = class_name
        p_cn.font.size = Pt(11)
        p_cn.font.bold = True
        p_cn.font.color.rgb = COLOR_WHITE
        p_cn.alignment = PP_ALIGN.CENTER

        # 2. Attributes Compartment
        attr_top = top + h_height + 0.05
        attr_height = 0.22 * len(attributes) + 0.15
        tb_a = slide.shapes.add_textbox(Inches(left + 0.1), Inches(attr_top), Inches(width - 0.2), Inches(attr_height))
        tf_a = tb_a.text_frame
        tf_a.word_wrap = True
        for i, a in enumerate(attributes):
            p = tf_a.paragraphs[0] if i == 0 else tf_a.add_paragraph()
            p.text = a
            p.font.size = Pt(8.5)
            p.font.name = "Consolas"
            p.font.color.rgb = COLOR_NAVY
            p.space_before = Pt(2) if i > 0 else Pt(0)

        # Divider Line between Attributes and Methods
        div_top = attr_top + attr_height
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.08), Inches(div_top), Inches(width - 0.16), Inches(0.015))
        div.fill.solid()
        div.fill.fore_color.rgb = COLOR_CARD_BORDER
        div.line.fill.background()

        # 3. Methods / Operations Compartment
        meth_top = div_top + 0.05
        meth_height = height - (meth_top - top) - 0.08
        tb_m = slide.shapes.add_textbox(Inches(left + 0.1), Inches(meth_top), Inches(width - 0.2), Inches(meth_height))
        tf_m = tb_m.text_frame
        tf_m.word_wrap = True
        for i, m in enumerate(methods):
            p = tf_m.paragraphs[0] if i == 0 else tf_m.add_paragraph()
            p.text = m
            p.font.size = Pt(8.5)
            p.font.name = "Consolas"
            p.font.color.rgb = COLOR_DARK_BLUE
            p.space_before = Pt(2) if i > 0 else Pt(0)

    # Class 1: Category
    draw_uml_class(s4, 0.8, 1.4, 3.4, 2.35, "Category", "Entity",
                   ["- id: int [PK]", "- name: string", "- description: string"],
                   ["+ getItems(): List<InventoryItem>", "+ getStorageProfile(): string"],
                   COLOR_DARK_BLUE)

    # Class 2: InventoryItem (Central)
    draw_uml_class(s4, 4.8, 1.4, 3.8, 3.1, "InventoryItem", "Aggregate Root",
                   ["- id: int [PK]", "- categoryId: int [FK]", "- name: string", "- unit: string", "- minThreshold: float"],
                   ["+ calculateCurrentStock(): float", "+ getNearestExpiration(): Date", "+ evaluateStatus(): StatusBadge", "+ isLowStock(): boolean"],
                   COLOR_NAVY)

    # Class 3: InventoryTransaction
    draw_uml_class(s4, 9.2, 1.4, 3.4, 3.1, "InventoryTransaction", "Entity",
                   ["- id: int [PK]", "- itemId: int [FK]", "- transactionType: 'IN'|'OUT'", "- quantity: float", "- purchaseDate: Date", "- expirationDate: Date", "- reason: string"],
                   ["+ isExpired(today: Date): boolean", "+ isExpiringSoon(days=3): boolean"],
                   COLOR_PRIMARY)

    # Class 4: AlertEngine (Service)
    draw_uml_class(s4, 0.8, 4.2, 3.4, 2.5, "AlertEngine", "Domain Service",
                   ["- alertHorizonDays: int = 3", "- replenishmentMultiplier: float = 2.0"],
                   ["+ evaluateExpiringLots(): List<Alert>", "+ evaluateShoppingList(): List<Item>", "+ calculateReorderQty(): float"],
                   COLOR_EMERALD)

    # Class 5: StockInDTO / StockOutDTO (Payload Data Structures)
    draw_uml_class(s4, 4.8, 4.75, 7.8, 1.95, "TransactionPayloads (DTOs)", "Data Transfer Objects",
                   ["+ StockInDTO: { itemId: int, quantity: float, purchaseDate: Date, expirationDate: Date, reason: string }",
                    "+ StockOutDTO: { itemId: int, quantity: float, reason: 'Used/Consumed' | 'Spoiled/Expired' | 'Transferred' }"],
                   ["+ validateInbound(): ValidationResult", "+ validateOutbound(currentStock: float): ValidationResult"],
                   RGBColor(71, 85, 105))

    # Relationship Connectors & Annotations
    # Connector 1: Category 1 --- 0..* InventoryItem
    tb_c1 = s4.shapes.add_textbox(Inches(4.15), Inches(2.2), Inches(0.7), Inches(0.4))
    p_c1 = tb_c1.text_frame.paragraphs[0]
    p_c1.text = "1 ─── 0..*\n«has»"
    p_c1.font.size = Pt(8)
    p_c1.font.bold = True
    p_c1.font.color.rgb = COLOR_PRIMARY
    p_c1.alignment = PP_ALIGN.CENTER

    # Connector 2: InventoryItem 1 --- 0..* InventoryTransaction
    tb_c2 = s4.shapes.add_textbox(Inches(8.55), Inches(2.2), Inches(0.7), Inches(0.4))
    p_c2 = tb_c2.text_frame.paragraphs[0]
    p_c2.text = "1 ─── 0..*\n«logs»"
    p_c2.font.size = Pt(8)
    p_c2.font.bold = True
    p_c2.font.color.rgb = COLOR_PRIMARY
    p_c2.alignment = PP_ALIGN.CENTER

    add_footer(s4, 4)

    # ==========================================
    # SLIDE 5: CLASS-RESPONSIBILITY-COLLABORATOR (CRC) CARDS
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "4. Class-Responsibility-Collaborator (CRC) Cards")

    def draw_crc_card(slide, left, top, width, height, class_name, stereotype, responsibilities, collaborators, header_color):
        # Outer index card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = header_color
        card.line.width = Pt(1.5)

        # Header bar
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.55))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = header_color
        hdr.line.fill.background()
        tf_h = hdr.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = f"CLASS: {class_name.upper()}   «{stereotype}»"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.alignment = PP_ALIGN.CENTER

        # Vertical Divider (60% Responsibilities / 40% Collaborators)
        resp_width = width * 0.60
        collab_width = width * 0.40 - 0.2
        collab_left = left + resp_width + 0.1

        v_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + resp_width), Inches(top + 0.62), Inches(0.015), Inches(height - 0.72))
        v_line.fill.solid()
        v_line.fill.fore_color.rgb = COLOR_CARD_BORDER
        v_line.line.fill.background()

        # Left Column: Responsibilities
        tb_r = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.62), Inches(resp_width - 0.25), Inches(height - 0.75))
        tf_r = tb_r.text_frame
        tf_r.word_wrap = True
        p_rt = tf_r.paragraphs[0]
        p_rt.text = "RESPONSIBILITIES (What it Knows & Does):"
        p_rt.font.size = Pt(9.5)
        p_rt.font.bold = True
        p_rt.font.color.rgb = header_color
        
        for r in responsibilities:
            pr = tf_r.add_paragraph()
            pr.text = f"• {r}"
            pr.font.size = Pt(9)
            pr.font.color.rgb = COLOR_TEXT_MAIN
            pr.space_before = Pt(4)

        # Right Column: Collaborators
        tb_c = slide.shapes.add_textbox(Inches(collab_left), Inches(top + 0.62), Inches(collab_width), Inches(height - 0.75))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        p_ct = tf_c.paragraphs[0]
        p_ct.text = "COLLABORATORS:"
        p_ct.font.size = Pt(9.5)
        p_ct.font.bold = True
        p_ct.font.color.rgb = COLOR_NAVY
        
        for c in collaborators:
            pc = tf_c.add_paragraph()
            pc.text = f"• {c}"
            pc.font.size = Pt(9)
            pc.font.color.rgb = COLOR_TEXT_MAIN
            pc.space_before = Pt(4)

    crc_list = [
        ("InventoryItem", "Aggregate Root",
         ["Maintains ingredient identity, category FK, unit, and min threshold.",
          "Calculates live aggregate stock from transaction history (IN - OUT).",
          "Determines nearest expiration date across active positive batches.",
          "Evaluates visual health status badge (Sufficient, Low Stock, Expired)."],
         ["Category", "InventoryTransaction", "AlertEngine"],
         COLOR_NAVY),

        ("InventoryTransaction", "Entity",
         ["Records atomic stock intake ('IN') with batch purchase & expiry dates.",
          "Logs consumption/deduction ('OUT') with designated reason classification.",
          "Validates batch expiration status against current date (FEFO rule).",
          "Preserves immutable audit trail of all physical pantry movements."],
         ["InventoryItem"],
         COLOR_PRIMARY),

        ("AlertEngine", "Domain Service",
         ["Continuously evaluates batch expiration horizons (3-day lookahead).",
          "Emits prioritized expiration alerts for the top dashboard banner.",
          "Scans all inventory items for deficits below minimum threshold.",
          "Auto-generates Reorder Shopping List with replenishment quantities."],
         ["InventoryItem", "InventoryTransaction"],
         COLOR_EMERALD)
    ]

    for idx, (cname, stype, resps, collabs, col) in enumerate(crc_list):
        draw_crc_card(s5, 0.8, 1.4 + idx * 1.82, 11.733, 1.68, cname, stype, resps, collabs, col)

    add_footer(s5, 5)

    # ==========================================
    # SLIDE 6: INBOUND STOCK-IN WORKFLOW
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "5. Inbound Stock-In Workflow & Category Filter Feature")

    steps_in = [
        ("Step 1: Open Stock-In Modal", "User clicks '+ In' on an inventory row or '+ Stock In' from the dashboard."),
        ("Step 2: Instant Category & Keyword Filter", "User selects a Category (e.g., 'Dairy Products') or types in the live search box. Modal dynamically filters dropdown to only matching materials."),
        ("Step 3: Quantity & Batch Dates Entry", "User specifies intake quantity, purchase date, and batch expiration date."),
        ("Step 4: Reason Tagging", "User selects reason: 'Initial Stock / Regular Purchase', 'Supplier Restock', or 'Fresh Produce Restock'."),
        ("Step 5: Database Commit & Real-Time Sync", "POST /api/transactions/in logs the transaction, updates stock totals, and refreshes the UI dashboard instantly.")
    ]

    for idx, (stitle, sdesc) in enumerate(steps_in):
        sbox = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.4 + idx * 1.05), Inches(11.3), Inches(0.95))
        sbox.fill.solid()
        sbox.fill.fore_color.rgb = COLOR_WHITE
        sbox.line.color.rgb = COLOR_PRIMARY
        sbox.line.width = Pt(1.5)

        num_badge = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.15), Inches(1.52 + idx * 1.05), Inches(0.7), Inches(0.7))
        num_badge.fill.solid()
        num_badge.fill.fore_color.rgb = COLOR_PRIMARY
        num_badge.line.fill.background()
        tf_nb = num_badge.text_frame
        p_nb = tf_nb.paragraphs[0]
        p_nb.text = f"{idx+1}"
        p_nb.font.size = Pt(18)
        p_nb.font.bold = True
        p_nb.font.color.rgb = COLOR_WHITE
        p_nb.alignment = PP_ALIGN.CENTER

        tb_step = s6.shapes.add_textbox(Inches(2.0), Inches(1.45 + idx * 1.05), Inches(10.1), Inches(0.85))
        tf_s = tb_step.text_frame
        tf_s.word_wrap = True
        p_st = tf_s.paragraphs[0]
        p_st.text = stitle
        p_st.font.size = Pt(11)
        p_st.font.bold = True
        p_st.font.color.rgb = COLOR_PRIMARY
        
        p_sd = tf_s.add_paragraph()
        p_sd.text = sdesc
        p_sd.font.size = Pt(9.5)
        p_sd.font.color.rgb = COLOR_TEXT_MAIN
        p_sd.space_before = Pt(2)
    add_footer(s6, 6)

    # ==========================================
    # SLIDE 7: OUTBOUND STOCK-OUT WORKFLOW
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "6. Outbound Stock-Out Workflow & Validation Rules")

    create_card(s7, 0.8, 1.4, 5.7, 5.35, "Stock-Out Execution Steps", COLOR_ROSE)
    tb_so = s7.shapes.add_textbox(Inches(0.95), Inches(1.95), Inches(5.4), Inches(4.6))
    tf_so = tb_so.text_frame
    tf_so.word_wrap = True
    so_steps = [
        ("1. Material Selection", "User filters by category or searches item name inside the Stock-Out modal."),
        ("2. Stock Sufficiency Check", "Modal validates that deduction quantity <= available stock balance. Over-deductions are blocked with an error banner."),
        ("3. Reason Classification", "User selects one of 3 clean reasons: 'Used / Consumed', 'Spoiled / Expired', or 'Transferred'."),
        ("4. Atomic Deduction Commit", "POST /api/transactions/out writes the 'OUT' transaction and updates the item's balance.")
    ]
    for idx, (sname, sdesc) in enumerate(so_steps):
        p = tf_so.paragraphs[0] if idx == 0 else tf_so.add_paragraph()
        p.text = sname
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_ROSE
        p.space_before = Pt(8) if idx > 0 else Pt(0)
        
        pd = tf_so.add_paragraph()
        pd.text = sdesc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = COLOR_TEXT_MAIN
        pd.space_before = Pt(2)

    create_card(s7, 6.8, 1.4, 5.7, 5.35, "Deduction Reason Definitions", COLOR_DARK_BLUE)
    tb_reasons = s7.shapes.add_textbox(Inches(6.95), Inches(1.95), Inches(5.4), Inches(4.6))
    tf_rea = tb_reasons.text_frame
    tf_rea.word_wrap = True
    reasons = [
        ("Used / Consumed (Default)", "Standard culinary usage: ingredients cooked, eaten, or prepared for meals."),
        ("Spoiled / Expired", "Discarded inventory: food that reached past expiration, spoiled, or went bad."),
        ("Transferred", "Relocated inventory: stock transferred to another kitchen, station, or storage area.")
    ]
    for idx, (rname, rdesc) in enumerate(reasons):
        p = tf_rea.paragraphs[0] if idx == 0 else tf_rea.add_paragraph()
        p.text = f"• {rname}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK_BLUE
        p.space_before = Pt(12) if idx > 0 else Pt(0)
        
        pd = tf_rea.add_paragraph()
        pd.text = rdesc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = COLOR_TEXT_MAIN
        pd.space_before = Pt(2)
    add_footer(s7, 7)

    # ==========================================
    # SLIDE 8: REAL-TIME HEALTH STATUS BADGES
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "7. Real-Time Telemetry & 5 Status Badges")

    badges_s8 = [
        ("SUFFICIENT (Green Badge)", "Trigger: current_stock >= min_threshold & No Expiring Batches", "Healthy stock level; no immediate action required.", COLOR_EMERALD),
        ("LOW_STOCK (Yellow Badge)", "Trigger: 0 < current_stock < min_threshold", "Inventory running low; automatically added to Reorder Shopping List.", RGBColor(234, 179, 8)),
        ("EXPIRING_SOON (Orange Badge)", "Trigger: 0 < (expiration_date - TODAY) <= 3 Days", "Batch near expiration; highlighted in top banner to prioritize usage.", COLOR_AMBER),
        ("EXPIRED (Red Badge)", "Trigger: Positive batch with expiration_date < TODAY", "Batch past expiration; flagged for immediate inspection or discard.", COLOR_ROSE),
        ("OUT_OF_STOCK (Red Badge)", "Trigger: current_stock <= 0", "Completely depleted; urgent restock required.", COLOR_NAVY)
    ]

    for idx, (bname, btrig, bact, bcol) in enumerate(badges_s8):
        bcard = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.4 + idx * 1.05), Inches(11.3), Inches(0.95))
        bcard.fill.solid()
        bcard.fill.fore_color.rgb = COLOR_WHITE
        bcard.line.color.rgb = bcol
        bcard.line.width = Pt(1.5)

        cbadge = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.15), Inches(1.55 + idx * 1.05), Inches(0.35), Inches(0.65))
        cbadge.fill.solid()
        cbadge.fill.fore_color.rgb = bcol
        cbadge.line.fill.background()

        tb_b = s8.shapes.add_textbox(Inches(1.65), Inches(1.45 + idx * 1.05), Inches(10.5), Inches(0.85))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True

        p_bt = tf_b.paragraphs[0]
        p_bt.text = bname
        p_bt.font.size = Pt(11)
        p_bt.font.bold = True
        p_bt.font.color.rgb = bcol

        p_bd = tf_b.add_paragraph()
        p_bd.text = f"{btrig}   •   {bact}"
        p_bd.font.size = Pt(9.5)
        p_bd.font.color.rgb = COLOR_TEXT_MAIN
        p_bd.space_before = Pt(2)
    add_footer(s8, 8)

    # ==========================================
    # SLIDE 9: REORDER SHOPPING LIST ENGINE
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "8. Automated Reorder Shopping List Engine")

    create_card(s9, 0.8, 1.4, 5.7, 5.35, "Reorder Calculation Formula", COLOR_PRIMARY)
    tb_calc = s9.shapes.add_textbox(Inches(0.95), Inches(1.95), Inches(5.4), Inches(4.6))
    tf_calc = tb_calc.text_frame
    tf_calc.word_wrap = True
    
    p_eq = tf_calc.paragraphs[0]
    p_eq.text = "Replenishment Target Formula:"
    p_eq.font.size = Pt(11)
    p_eq.font.bold = True
    p_eq.font.color.rgb = COLOR_PRIMARY
    
    p_f1 = tf_calc.add_paragraph()
    p_f1.text = "Target Stock Level = 2 × min_threshold\nSuggested Reorder = MAX(1, Target - Current Stock)"
    p_f1.font.size = Pt(10)
    p_f1.font.name = "Consolas"
    p_f1.font.color.rgb = COLOR_NAVY
    p_f1.space_before = Pt(4)

    p_ex = tf_calc.add_paragraph()
    p_ex.text = "\nExample Calculation:"
    p_ex.font.size = Pt(11)
    p_ex.font.bold = True
    p_ex.font.color.rgb = COLOR_PRIMARY

    p_exd = tf_calc.add_paragraph()
    p_exd.text = "• Item: Lemons (Min: 10 units, Current: 2 units)\n• Target: 2 × 10 = 20 units\n• Suggested Reorder: 20 - 2 = 18 units (Urgency: HIGH)\n\n• Item: Strawberries (Min: 1.5 kg, Current: 0 kg)\n• Target: 2 × 1.5 = 3.0 kg\n• Suggested Reorder: 3.0 kg (Urgency: CRITICAL)"
    p_exd.font.size = Pt(9.5)
    p_exd.font.color.rgb = COLOR_TEXT_MAIN
    p_exd.space_before = Pt(3)

    create_card(s9, 6.8, 1.4, 5.7, 5.35, "Shopping List UI & 1-Click Quick Restock", COLOR_EMERALD)
    tb_sl = s9.shapes.add_textbox(Inches(6.95), Inches(1.95), Inches(5.4), Inches(4.6))
    tf_sl = tb_sl.text_frame
    tf_sl.word_wrap = True
    sl_features = [
        ("⚡ 1-Click Instant Restock", "Click '+[Qty]' directly in modal to stock in suggested quantity without leaving the view."),
        ("Real-Time Red Badge Counter", "Header badge automatically tracks the exact count of items requiring restocking."),
        ("Search Inside Modal", "Instant search bar allows quick filtering of large shopping lists by name or category."),
        ("Print / Export Support", "One-click 'Print / Export List' formats the table cleanly for mobile or physical grocery shopping.")
    ]
    for idx, (ftitle, fdesc) in enumerate(sl_features):
        p = tf_sl.paragraphs[0] if idx == 0 else tf_sl.add_paragraph()
        p.text = f"✔ {ftitle}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_EMERALD
        p.space_before = Pt(8) if idx > 0 else Pt(0)
        
        pd = tf_sl.add_paragraph()
        pd.text = fdesc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = COLOR_TEXT_MAIN
        pd.space_before = Pt(2)
    add_footer(s9, 9)

    # ==========================================
    # SLIDE 10: CORE MODULES & CODE IMPLEMENTATION
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "9. Core Modules & Code Implementation Breakdown")

    create_card(s10, 0.8, 1.4, 5.7, 5.35, "Backend Modules (Node / Express / SQLite)", COLOR_PRIMARY)
    tb_m10 = s10.shapes.add_textbox(Inches(0.95), Inches(1.95), Inches(5.4), Inches(4.6))
    tf_m10 = tb_m10.text_frame
    tf_m10.word_wrap = True
    be_modules = [
        ("server.js", "Express REST server hosting /api/categories, /api/items, /api/transactions/in, /api/transactions/out, /api/dashboard/summary, and /api/alerts/*."),
        ("db.js", "SQLite database connection manager with async run(), get(), all() helper wrappers and auto-initialization."),
        ("schema.sql & seed.sql", "DDL script defining categories, items, and inventory_transactions with realistic sample items.")
    ]
    for idx, (mtitle, mdesc) in enumerate(be_modules):
        p = tf_m10.paragraphs[0] if idx == 0 else tf_m10.add_paragraph()
        p.text = f"• {mtitle}"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_before = Pt(6) if idx > 0 else Pt(0)
        
        pd = tf_m10.add_paragraph()
        pd.text = mdesc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = COLOR_TEXT_MAIN
        pd.space_before = Pt(2)

    create_card(s10, 6.8, 1.4, 5.7, 5.35, "Frontend Modules (React 18 + Vite)", COLOR_DARK_BLUE)
    tb_fe10 = s10.shapes.add_textbox(Inches(6.95), Inches(1.95), Inches(5.4), Inches(4.6))
    tf_fe10 = tb_fe10.text_frame
    tf_fe10.word_wrap = True
    fe_modules = [
        ("App.jsx", "Root orchestrator managing live data fetch, search/category state, alert banners, and modal sync."),
        ("StockInModal.jsx / StockOutModal.jsx", "Modals featuring live category dropdown filtering, keyword search, and validation."),
        ("InventoryTable.jsx & CategoryGrid.jsx", "Renderable interactive tables and cards with color-coded status badges."),
        ("ShoppingListModal.jsx", "Reorder modal displaying low-stock items with 1-click Quick Restock and export tools.")
    ]
    for idx, (vtitle, vdesc) in enumerate(fe_modules):
        p = tf_fe10.paragraphs[0] if idx == 0 else tf_fe10.add_paragraph()
        p.text = f"• {vtitle}"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK_BLUE
        p.space_before = Pt(6) if idx > 0 else Pt(0)
        
        pd = tf_fe10.add_paragraph()
        pd.text = vdesc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = COLOR_TEXT_MAIN
        pd.space_before = Pt(2)
    add_footer(s10, 10)

    # ==========================================
    # SLIDE 11: SUMMARY & LIVE VERCEL DEPLOYMENT
    # ==========================================
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "10. Summary & Live Prototype Deployment")

    summary_cards = [
        ("1. Accurate Stock Tracking", "Maintains exact real-time quantities and batch expiration dates across all 6 food groups with zero data loss.", COLOR_PRIMARY),
        ("2. Instant Material Discovery", "Category filter dropdowns and keyword search in Stock-In/Out modals allow users to find items in seconds.", COLOR_EMERALD),
        ("3. Proactive Spoilage Prevention", "3-day lookahead expiration warnings and nearest expiry tags ensure perishable ingredients are used before spoiling.", COLOR_AMBER),
        ("4. Automated Reordering", "Calculates required restock amounts automatically when stock falls below threshold, simplifying grocery planning.", COLOR_DARK_BLUE)
    ]

    for idx, (stitle, sdesc, scolor) in enumerate(summary_cards):
        r = idx // 2
        c = idx % 2
        card = create_card(s11, 1.0 + c * 5.8, 1.38 + r * 2.2, 5.5, 2.05, stitle, scolor)
        tb = s11.shapes.add_textbox(Inches(1.2 + c * 5.8), Inches(1.92 + r * 2.2), Inches(5.1), Inches(1.35))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sdesc
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.line_spacing = 1.25

    # Live Vercel Deployment Link Banner
    live_banner = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.95), Inches(11.3), Inches(0.95))
    live_banner.fill.solid()
    live_banner.fill.fore_color.rgb = COLOR_NAVY
    live_banner.line.color.rgb = COLOR_EMERALD
    live_banner.line.width = Pt(2)

    tb_live = s11.shapes.add_textbox(Inches(1.2), Inches(6.0), Inches(10.9), Inches(0.85))
    tf_l = tb_live.text_frame
    tf_l.word_wrap = True
    
    p_l1 = tf_l.paragraphs[0]
    p_l1.text = "🌐 LIVE PRODUCTION DEPLOYMENT (VERCEL)"
    p_l1.font.size = Pt(10)
    p_l1.font.bold = True
    p_l1.font.color.rgb = COLOR_ACCENT
    
    p_l2 = tf_l.add_paragraph()
    p_l2.text = "https://refrigerator-inventory-system.vercel.app/"
    p_l2.font.size = Pt(15)
    p_l2.font.bold = True
    p_l2.font.color.rgb = COLOR_EMERALD
    p_l2.space_before = Pt(2)

    add_footer(s11, 11)

    # Save presentation
    output_files = ['RIMS_System_Specification_Updated.pptx', 'RIMS_System_Specification.pptx']
    for out_path in output_files:
        try:
            prs.save(out_path)
            print(f"Successfully saved {TOTAL_SLIDES}-slide presentation to: {out_path}")
        except PermissionError:
            print(f"Note: Could not overwrite '{out_path}' because it is currently open in PowerPoint.")

if __name__ == "__main__":
    build_presentation()
