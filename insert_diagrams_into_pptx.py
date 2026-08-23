import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def update_user_presentation():
    pptx_path = "RIMS_System_Specification.pptx"
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return

    try:
        prs = Presentation(pptx_path)
    except PermissionError:
        print("Please close PowerPoint if 'RIMS_System_Specification.pptx' is currently open.")
        return

    print(f"Loaded existing presentation with {len(prs.slides)} slides.")

    # Palette
    COLOR_NAVY = RGBColor(15, 23, 42)
    COLOR_ACCENT = RGBColor(14, 165, 233)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139)

    def add_header(slide, title_text):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_NAVY
        bar.line.fill.background()

        tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.5), Inches(0.3))
        tf_cat = tb_cat.text_frame
        p_c = tf_cat.paragraphs[0]
        p_c.text = "PROTOTYPE SYSTEM SPECIFICATION"
        p_c.font.size = Pt(9.5)
        p_c.font.bold = True
        p_c.font.color.rgb = COLOR_ACCENT

        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(11.5), Inches(0.65))
        tf_title = tb_title.text_frame
        p_t = tf_title.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE

    def add_footer(slide, slide_num, total_slides):
        tb_f = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.35))
        tf_f = tb_f.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"Refrigerator & Pantry Inventory Subsystem (RIMS)   |   Slide {slide_num} of {total_slides}"
        p_f.font.size = Pt(9)
        p_f.font.color.rgb = COLOR_TEXT_MUTED

    # If slides 4 and 5 exist, update them with the high-res diagrams
    total = len(prs.slides)
    if total >= 4 and os.path.exists("uml_class_diagram.png"):
        s4 = prs.slides[3] # Slide 4 (0-indexed 3)
        # Clear shapes on slide 4 except background
        for shape in list(s4.shapes):
            sp = shape._element
            sp.getparent().remove(sp)
        add_header(s4, "3. Object-Oriented UML Class Diagram")
        s4.shapes.add_picture("uml_class_diagram.png", Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.6))
        add_footer(s4, 4, total)
        print("Updated Slide 4 with high-res UML Class Diagram.")

    if total >= 5 and os.path.exists("crc_cards_diagram.png"):
        s5 = prs.slides[4] # Slide 5 (0-indexed 4)
        for shape in list(s5.shapes):
            sp = shape._element
            sp.getparent().remove(sp)
        add_header(s5, "4. Class-Responsibility-Collaborator (CRC) Cards")
        s5.shapes.add_picture("crc_cards_diagram.png", Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.6))
        add_footer(s5, 5, total)
        print("Updated Slide 5 with high-res CRC Cards Diagram.")

    try:
        prs.save("RIMS_System_Specification.pptx")
        print("Successfully saved updated 'RIMS_System_Specification.pptx'!")
    except PermissionError:
        prs.save("RIMS_System_Specification_With_Diagrams.pptx")
        print("PowerPoint is open, saved copy as 'RIMS_System_Specification_With_Diagrams.pptx'!")

if __name__ == "__main__":
    update_user_presentation()
